# generators/generateur_pptx.py
"""
Générateur PowerPoint Orange Business — 9 slides
Charte graphique officielle Orange : #FF7900
Structure identique au modèle Dinootoo PPTX Generator
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Couleurs officielles Orange Business ────────────────────────────────────
ORANGE = RGBColor(0xFF, 0x79, 0x00)
BLANC  = RGBColor(0xFF, 0xFF, 0xFF)
NOIR   = RGBColor(0x00, 0x00, 0x00)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
GRIS_L = RGBColor(0xF2, 0xF2, 0xF2)
SAUMON = RGBColor(0xFF, 0xCC, 0xAA)

W = Inches(13.33)
H = Inches(7.5)

# ── Chemins assets ────────────────────────────────────────────────────────────
_BASE         = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS_ICONS  = os.path.join(_BASE, "Assets", "icons")
ASSETS_IMAGES = os.path.join(_BASE, "Assets", "images")
LOGO_ORANGE   = os.path.join(ASSETS_IMAGES, "orange logo.png")


# ═══════════════════════════════════════════════════════════════════════════
# POINT D'ENTRÉE
# ═══════════════════════════════════════════════════════════════════════════

def generer_pptx(data: dict, textes: dict, chemin: str):
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    _slide_couverture(prs, data)              # Slide 1
    _slide_pourquoi(prs, data, textes)        # Slide 2
    _slide_benefices(prs, data, textes)       # Slide 3
    _slide_fonctionnalites(prs, data)         # Slide 4
    _slide_cas_usage(prs, data, textes)       # Slide 5
    _slide_securite(prs, data, textes)        # Slide 6
    _slide_accompagnement(prs, data, textes)  # Slide 7
    _slide_tarification(prs, data)            # Slide 8
    _slide_merci(prs, data, textes)           # Slide 9

    prs.save(chemin)
    print(f"  [PPTX] Sauvegarde : {chemin}")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COUVERTURE
# ═══════════════════════════════════════════════════════════════════════════

def _slide_couverture(prs, d):
    s = _blank(prs)

    # Date — haut droite en orange
    _txt(s, Inches(9.5), Inches(0.3), Inches(3.5), Inches(0.4),
         d["date"], Pt(12), ORANGE, bold=True, align=PP_ALIGN.RIGHT)

    # Titre principal orange (grand, gras)
    _txt(s, Inches(0.5), Inches(1.0), Inches(8.5), Inches(2.2),
         f"Offre commerciale\n{d['titre_offre']}", Pt(40), ORANGE, bold=True)

    # Sous-titre noir gras
    _txt(s, Inches(0.5), Inches(3.3), Inches(8.5), Inches(0.9),
         d["sous_titre"], Pt(18), NOIR, bold=True)

    # Barre noire en bas
    _rect(s, 0, H - Inches(0.6), W, Inches(0.6), DARK)

    # Logo Orange officiel bas gauche
    if os.path.exists(LOGO_ORANGE):
        s.shapes.add_picture(LOGO_ORANGE,
                             Inches(0.3), H - Inches(1.3),
                             Inches(1.2), Inches(0.8))

    # Mention IA bas droite dans la barre noire
    _txt(s, Inches(7.0), H - Inches(0.55), Inches(6.0), Inches(0.45),
         "Généré par IA avec l'outil Orange Business AI",
         Pt(8), BLANC, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — POURQUOI CETTE OFFRE ?
# ═══════════════════════════════════════════════════════════════════════════

def _slide_pourquoi(prs, d, textes):
    s = _blank(prs)
    _footer_restricted(s, 2)

    # Grand titre noir — colonne gauche
    _txt(s, Inches(0.5), Inches(0.8), Inches(7.5), Inches(2.0),
         f"Pourquoi {d['titre_offre']} ?", Pt(36), NOIR, bold=True)

    # Texte LLM sous le titre
    _txt(s, Inches(0.5), Inches(3.0), Inches(7.2), Inches(3.0),
         textes.get("pourquoi", ""), Pt(13), NOIR, bold=True)

    # Illustration droite
    _img(s, "Shared goals-rafiki.png",
         Inches(7.8), Inches(0.5), Inches(5.2), Inches(6.0))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — BÉNÉFICES
# ═══════════════════════════════════════════════════════════════════════════

def _slide_benefices(prs, d, textes):
    s = _blank(prs)
    _header_bande(s, 3)

    # Titre orange
    _txt(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
         "Les bénéfices pour votre entreprise", Pt(15), ORANGE, bold=True)
    # Sous-titre noir
    _txt(s, Inches(0.5), Inches(1.8), Inches(12), Inches(0.5),
         "Productivité, mobilité et collaboration", Pt(15), NOIR, bold=True)

    # 3 paragraphes bénéfices
    y = Inches(2.6)
    for key in ["benefice_1", "benefice_2", "benefice_3"]:
        _txt(s, Inches(0.5), y, Inches(12.3), Inches(0.75),
             textes.get(key, ""), Pt(12), NOIR, bold=True)
        y += Inches(1.0)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — FONCTIONNALITÉS CLÉS
# ═══════════════════════════════════════════════════════════════════════════

def _slide_fonctionnalites(prs, d):
    s = _blank(prs)
    _header_bande(s, 4)

    # Titre orange + sous-titre noir
    _txt(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
         "Fonctionnalités clés de votre solution", Pt(15), ORANGE, bold=True)
    _txt(s, Inches(0.5), Inches(1.8), Inches(12), Inches(0.5),
         "Des outils puissants pour votre quotidien", Pt(14), NOIR, bold=True)

    boites  = _construire_boites(d)
    icones  = [
        os.path.join(ASSETS_ICONS, "collaboration.png"),
        os.path.join(ASSETS_ICONS, "smartphone.png"),
        os.path.join(ASSETS_ICONS, "meeting.png"),
    ]
    x_pos = [Inches(0.5), Inches(4.6), Inches(8.7)]
    box_w  = Inches(3.8)
    box_h  = Inches(1.6)

    for i, (titre, texte) in enumerate(boites[:3]):
        x = x_pos[i]
        # Boîte orange
        _rect(s, x, Inches(2.5), box_w, box_h, ORANGE)
        # Icône centrée dans la boîte
        ic = icones[i]
        if os.path.exists(ic):
            s.shapes.add_picture(ic,
                                    x + Inches(1.55), Inches(2.6),
                                    Inches(0.7), Inches(0.7))
        # Titre boîte
        _txt(s, x, Inches(2.65), box_w, Inches(0.8),
                titre, Pt(12), BLANC, bold=True, align=PP_ALIGN.CENTER)
        # Texte description sous la boîte
        _txt(s, x, Inches(4.2), box_w, Inches(2.0),
                texte, Pt(10), NOIR)


def _construire_boites(d):
    if d.get("a_fibre") and d.get("a_microsoft"):
        return [
            ("Collaboration en temps réel",
                "Travaillez ensemble sur vos documents avec OneDrive et SharePoint, "
                "même à distance."),
            ("Mobilité et accès sécurisé",
                f"Fibre {d.get('fibre_debit','')} Mbps garantis. Connexion dédiée "
                "FTTO pour votre entreprise."),
            ("Productivité automatisée",
                "Pack Office complet, Teams et outils Microsoft 365 pour gagner "
                "du temps au quotidien."),
        ]
    elif d.get("a_fibre"):
        return [
            ("Débit Garanti",
                f"{d.get('fibre_debit','')} Mbps dédiés. Connexion symétrique "
                "sans partage de bande passante."),
            ("Fiabilité Maximale",
                "Engagement sur la qualité de service. Support technique "
                "Orange Business 24h/24."),
            ("Installation Complète",
                "Travaux FO, routeur et mise en service pris en charge "
                "par les équipes Orange Business."),
        ]
    else:
        ms = d.get("microsoft_plan", "Microsoft 365")
        return [
            ("Collaboration en temps réel",
                "Travaillez simultanément sur vos documents et échangez "
                "via Teams pour une collaboration fluide."),
            ("Mobilité et accès sécurisé",
                "Accédez à vos applications et fichiers depuis n'importe "
                "quel appareil, en toute sécurité."),
            ("Productivité automatisée",
                f"Suite {ms} : Word, Excel, PowerPoint et tous les outils "
                "pour votre quotidien."),
        ]


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CAS D'USAGE CLIENT
# ═══════════════════════════════════════════════════════════════════════════

def _slide_cas_usage(prs, d, textes):
    s = _blank(prs)
    _header_bande(s, 5)

    # Titre orange
    _txt(s, Inches(0.5), Inches(1.3), Inches(6.5), Inches(0.5),
            "Cas d'usage client", Pt(15), ORANGE, bold=True)

    # Texte gauche
    _txt(s, Inches(0.5), Inches(2.2), Inches(6.5), Inches(3.5),
            textes.get("cas_usage", ""), Pt(13), NOIR, bold=True)

    # Image droite — séparée par une ligne verticale légère
    _rect(s, Inches(7.3), Inches(0.9), Inches(0.03), Inches(5.8), GRIS_L)
    _img(s, "At the office-rafiki.png",
            Inches(7.5), Inches(0.9), Inches(5.5), Inches(5.8))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SÉCURITÉ ET CONFORMITÉ
# ═══════════════════════════════════════════════════════════════════════════

def _slide_securite(prs, d, textes):
    s = _blank(prs)
    _header_bande(s, 6)

    _txt(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
            "Sécurité et conformité", Pt(15), ORANGE, bold=True)
    _txt(s, Inches(0.5), Inches(1.8), Inches(12), Inches(0.5),
            "Vos données protégées au quotidien", Pt(14), NOIR, bold=True)

    securite = textes.get("securite", "")
    phrases  = [p.strip() for p in securite.split(".") if p.strip()]
    mid      = max(1, len(phrases) // 2)
    col_g    = ". ".join(phrases[:mid]) + "."
    col_d    = ". ".join(phrases[mid:]) + ("." if phrases[mid:] else "")

    _txt(s, Inches(0.5), Inches(2.6), Inches(5.8), Inches(3.5),
            col_g, Pt(12), NOIR, bold=True)
    if col_d.strip(".").strip():
        _txt(s, Inches(7.0), Inches(2.6), Inches(5.8), Inches(3.5),
                col_d, Pt(12), NOIR, bold=True)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ACCOMPAGNEMENT ET SUPPORT
# ═══════════════════════════════════════════════════════════════════════════

def _slide_accompagnement(prs, d, textes):
    s = _blank(prs)
    _header_bande(s, 7)

    _txt(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
            "Accompagnement et support", Pt(15), ORANGE, bold=True)
    _txt(s, Inches(0.5), Inches(1.8), Inches(12), Inches(0.5),
            "Un partenaire à vos côtés", Pt(14), NOIR, bold=True)

    y = Inches(2.7)
    for key in ["accompagnement_1", "accompagnement_2", "accompagnement_3"]:
        txt = textes.get(key, "")
        if txt:
            _txt(s, Inches(0.5), y, Inches(12.3), Inches(0.75),
                    txt, Pt(12), NOIR, bold=True)
            y += Inches(1.0)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TARIFICATION ET DEVIS
# ═══════════════════════════════════════════════════════════════════════════

def _slide_tarification(prs, d):
    s = _blank(prs)
    _header_bande(s, 8)

    _txt(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
            "Tarification et devis", Pt(15), ORANGE, bold=True)
    _txt(s, Inches(0.5), Inches(1.8), Inches(12), Inches(0.5),
            "Détail des licences et coûts", Pt(14), NOIR, bold=True)

    headers = ["Licence / Offre", "Quantité", "Prix unitaire", "Prix total (TND/an)"]
    lignes  = _lignes_tarif(d)
    _tableau_tarif(s, headers, lignes, Inches(0.5), Inches(2.5), Inches(12.3))


def _lignes_tarif(d) -> list:
    lignes = []
    fibre = d.get("fibre")
    if fibre:
        lignes.append([
            fibre.get("nom_offre", "Fibre Orange Business"),
            "1",
            f"{fibre.get('prix_mensuel_total', 0):.2f} TND/mois",
            f"{fibre.get('prix_annuel', 0):,.2f}",
        ])
    ms = d.get("microsoft")
    if ms:
        lignes.append([
            ms.get("nom_produit", "Microsoft 365"),
            str(ms.get("nombre_licences", 0)),
            f"{ms.get('prix_unitaire_tnd', 0):.2f} TND/an/licence",
            f"{ms.get('prix_annuel', 0):,.2f}",
        ])
    lignes.append(["TOTAL ANNUEL", "", "", f"{d.get('prix_total_annuel', 0):,.2f}"])
    return lignes


def _tableau_tarif(slide, headers, rows, left, top, width):
    nb_r  = len(rows) + 1
    nb_c  = len(headers)
    row_h = Inches(0.6)

    tbl = slide.shapes.add_table(nb_r, nb_c, left, top, width, row_h * nb_r).table
    col_widths = [int(width * p) for p in [0.40, 0.14, 0.24, 0.22]]
    for j, cw in enumerate(col_widths):
        tbl.columns[j].width = cw

    # En-tête orange
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ORANGE
        p   = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.name  = "Helvetica Neue"
        run.font.bold  = True
        run.font.size  = Pt(11)
        run.font.color.rgb = BLANC

    # Lignes données
    for i, row in enumerate(rows, start=1):
        est_total = (i == len(rows))
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            if est_total:
                cell.fill.fore_color.rgb = ORANGE
            elif i % 2 == 0:
                cell.fill.fore_color.rgb = SAUMON
            else:
                cell.fill.fore_color.rgb = BLANC
            p   = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            run.font.name  = "Helvetica Neue"
            run.font.size  = Pt(11)
            run.font.bold  = est_total or (j == 0)
            run.font.color.rgb = BLANC if est_total else NOIR


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — MERCI
# ═══════════════════════════════════════════════════════════════════════════

def _slide_merci(prs, d, textes):
    s = _blank(prs)
    _footer_restricted(s, 9)

    # Grand titre orange
    _txt(s, Inches(0.5), Inches(1.2), Inches(7.5), Inches(1.8),
            "Merci !", Pt(54), ORANGE, bold=True)

    # Sous-titre noir gras
    _txt(s, Inches(0.5), Inches(3.0), Inches(7.5), Inches(0.8),
            "Contact et prochaines étapes", Pt(22), NOIR, bold=True)

    # Conclusion LLM
    _txt(s, Inches(0.5), Inches(4.0), Inches(7.5), Inches(1.5),
            textes.get("conclusion", ""), Pt(12), NOIR)

    # Illustration droite
    _img(s, "Company-rafiki.png",
            Inches(8.3), Inches(0.8), Inches(4.7), Inches(6.0))


# ═══════════════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════════════

def _blank(prs):
    """Slide vraiment vierge — supprime tous les placeholders."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Supprimer les placeholders résiduels ("Click to add...")
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)
    return slide


def _header_bande(slide, page_num):
    """Bande noire en haut + numéro de page + Orange Restricted haut et bas."""
    # Numéro de page haut gauche (au-dessus de la bande)
    _txt(slide, Inches(0.3), Inches(0.05), Inches(0.8), Inches(0.35),
            str(page_num), Pt(9), RGBColor(0x88, 0x88, 0x88))
    # "Orange Restricted" centré au-dessus de la bande
    _txt(slide, Inches(4.5), Inches(0.05), Inches(4.5), Inches(0.35),
            "Orange Restricted", Pt(9), ORANGE, align=PP_ALIGN.CENTER)
    # Bande noire
    _rect(slide, 0, Inches(0.45), W, Inches(0.35), DARK)
    # Footer bas
    _footer_restricted(slide, None)


def _footer_restricted(slide, page_num):
    """Footer 'Orange Restricted' centré en bas + numéro si fourni."""
    _txt(slide, Inches(4.5), H - Inches(0.35), Inches(4.5), Inches(0.3),
            "Orange Restricted", Pt(8), ORANGE, align=PP_ALIGN.CENTER)
    if page_num:
        _txt(slide, Inches(0.3), H - Inches(0.35), Inches(0.8), Inches(0.3),
                str(page_num), Pt(8), RGBColor(0x88, 0x88, 0x88))


def _rect(slide, left, top, w, h, color):
    shp = slide.shapes.add_shape(1, left, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()


def _txbox(slide, left, top, w, h):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tb.text_frame.word_wrap = True
    tb.line.fill.background()   # ← supprime la bordure visible
    return tb


def _txt(slide, left, top, w, h, texte, size, color,
            bold=False, italic=False, align=PP_ALIGN.LEFT):
    tb = _txbox(slide, left, top, w, h)
    p  = tb.text_frame.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text           = str(texte)
    r.font.name      = "Helvetica Neue"
    r.font.size      = size
    r.font.color.rgb = color
    r.font.bold      = bold
    r.font.italic    = italic


def _img(slide, nom, left, top, w, h):
    chemin = os.path.join(ASSETS_IMAGES, nom)
    if os.path.exists(chemin):
        slide.shapes.add_picture(chemin, left, top, w, h)
